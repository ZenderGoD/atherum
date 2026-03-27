"""Generate Atherum Stakeholder Deck — PPTX"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Colors
BG = RGBColor(0x0D, 0x0D, 0x14)
FG = RGBColor(0xE8, 0xE8, 0xEC)
DIM = RGBColor(0x88, 0x88, 0x99)
ACCENT = RGBColor(0x5A, 0xB0, 0x9E)  # muted teal
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_SURFACE = RGBColor(0x14, 0x14, 0x1E)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_slide():
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return slide

def add_text(slide, left, top, width, height, text, size=18, color=FG, bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tf

def add_bullet_list(slide, left, top, width, height, items, size=16, color=FG, bullet_color=ACCENT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(8)
        p.level = 0
    return tf

def section_header(slide, label, title):
    add_text(slide, 0.8, 0.6, 5, 0.5, label.upper(), size=12, color=ACCENT, bold=True)
    add_text(slide, 0.8, 1.0, 11, 1.0, title, size=36, color=WHITE, bold=True)

def slide_number(slide, num):
    add_text(slide, 12.3, 7.0, 0.8, 0.4, str(num), size=10, color=DIM, align=PP_ALIGN.RIGHT)

# ─── Slide 1: Title ──────────────────────────────────────────────────
s = add_slide()
add_text(s, 0.8, 2.0, 11, 1.2, "Atherum", size=60, color=WHITE, bold=True)
add_text(s, 0.8, 3.3, 11, 0.8, "Collective Intelligence Engine", size=28, color=ACCENT)
add_text(s, 0.8, 4.3, 11, 0.6, "10 AI agents deliberate on your content. One API call.", size=18, color=DIM)
add_text(s, 0.8, 6.5, 6, 0.4, "Zeus Ecosystem  |  Apache 2.0  |  github.com/ZenderGoD/atherum", size=11, color=DIM)
slide_number(s, 1)

# ─── Slide 2: The Problem ────────────────────────────────────────────
s = add_slide()
section_header(s, "Problem", "Single perspectives are not enough")
add_bullet_list(s, 0.8, 2.2, 11, 4.5, [
    "Single-model AI reviews give one perspective — biased, shallow, no debate",
    "Traditional focus groups cost $10K-$50K per session and take 2+ weeks",
    "Content teams publish without structured quality feedback",
    "Social listening is backward-looking — tells you what happened, not what will happen",
    "Market research is an $80B+ industry doing the same thing it did 20 years ago",
], size=18, color=FG)
slide_number(s, 2)

# ─── Slide 3: The Solution ───────────────────────────────────────────
s = add_slide()
section_header(s, "Solution", "Synthetic collective intelligence")
add_bullet_list(s, 0.8, 2.2, 11, 4.5, [
    "Multiple AI agents with distinct personas, expertise, and reasoning styles",
    "They debate through structured rounds and converge on collective verdicts",
    "Not a single opinion — a deliberated consensus with dissent tracking",
    "One API call  ->  10 perspectives  ->  structured verdict in under 60 seconds",
    "Follow-up questions answered with full deliberation context",
], size=18, color=FG)
slide_number(s, 3)

# ─── Slide 4: How It Works ───────────────────────────────────────────
s = add_slide()
section_header(s, "Process", "How it works")
steps = [
    ("1", "Submit", "POST content description + image URL via API"),
    ("2", "Analyze", "10 agents independently evaluate (Round 1, vision-enabled)"),
    ("3", "Deliberate", "Anonymous summaries shared, agents update positions (Rounds 2-3)"),
    ("4", "Converge", "TF-IDF similarity measures agreement across agents"),
    ("5", "Verdict", "Synthesis: approval score, agreements, dissent, minority report"),
    ("6", "Ask", "Follow-up questions with full deliberation transcript as context"),
]
for i, (num, title, desc) in enumerate(steps):
    y = 2.2 + i * 0.8
    add_text(s, 0.8, y, 0.5, 0.5, num, size=20, color=ACCENT, bold=True)
    add_text(s, 1.5, y, 2.0, 0.5, title, size=18, color=WHITE, bold=True)
    add_text(s, 3.8, y, 8.5, 0.5, desc, size=16, color=DIM)
slide_number(s, 4)

# ─── Slide 5: The Agents ─────────────────────────────────────────────
s = add_slide()
section_header(s, "Agents", "10 specialized reviewers")
agents = [
    ("Target Audience Member", "analytical"),
    ("Brand Critic", "creative"),
    ("Trend Analyst", "skeptical"),
    ("Marketing Expert", "pragmatic"),
    ("Social Media User", "synthesizing"),
    ("Creative Director", "visionary"),
    ("UX Designer", "analytical"),
    ("E-commerce Specialist", "creative"),
    ("Consumer Psychologist", "skeptical"),
    ("Photographer", "pragmatic"),
]
for i, (name, style) in enumerate(agents):
    col = i % 2
    row = i // 2
    x = 0.8 + col * 6.0
    y = 2.3 + row * 0.9
    add_text(s, x, y, 4.0, 0.4, name, size=16, color=WHITE, bold=True)
    add_text(s, x + 4.0, y, 2.0, 0.4, style, size=14, color=ACCENT)
slide_number(s, 5)

# ─── Slide 6: Real Output ────────────────────────────────────────────
s = add_slide()
section_header(s, "Output", "Real review result")
add_text(s, 0.8, 2.2, 5, 0.4, "Content: Minimalist sneaker on concrete, dramatic side lighting", size=14, color=DIM)
add_text(s, 0.8, 2.8, 2, 0.6, "16/100", size=48, color=ACCENT, bold=True)
add_text(s, 3.0, 2.95, 2, 0.4, "approval score", size=14, color=DIM)
add_text(s, 0.8, 3.8, 11, 0.8, '"Technically competent but strategically hollow. Generic urban minimalism fails to build brand resonance."', size=16, color=FG)
add_text(s, 0.8, 4.8, 5, 0.4, "Key Agreements:", size=14, color=WHITE, bold=True)
add_bullet_list(s, 0.8, 5.2, 11, 2.0, [
    "Technical quality is high",
    "Aesthetic is overused and generic",
    "Lacks distinct brand personality",
    "Functions as commodity, not compelling narrative",
], size=14, color=DIM)
add_text(s, 8.0, 2.8, 4, 0.4, "3 agents  |  2 rounds  |  <20 seconds", size=13, color=DIM)
slide_number(s, 6)

# ─── Slide 7: Four Engines ───────────────────────────────────────────
s = add_slide()
section_header(s, "Engines", "Four engines, one platform")
engines = [
    ("Mirage", "Multi-agent swarm deliberation", "BUILT"),
    ("OASIS", "Social media simulation (Twitter/Reddit)", "PLANNED"),
    ("Atlas", "Knowledge graph construction", "PLANNED"),
    ("Scribe", "Intelligent report generation", "PLANNED"),
]
for i, (name, desc, status) in enumerate(engines):
    y = 2.5 + i * 1.1
    add_text(s, 0.8, y, 3, 0.5, name, size=24, color=ACCENT, bold=True)
    add_text(s, 3.5, y + 0.05, 7, 0.5, desc, size=16, color=FG)
    sc = ACCENT if status == "BUILT" else DIM
    add_text(s, 11.0, y + 0.05, 1.5, 0.5, status, size=12, color=sc, bold=True)
slide_number(s, 7)

# ─── Slide 8: Platform Capabilities ──────────────────────────────────
s = add_slide()
section_header(s, "Platform", "What makes it different")
caps = [
    ("Mnemosyne", "Persistent agent memory across sessions"),
    ("Strata", "Tiered context loading (L0/L1/L2)"),
    ("Aegis", "Cost governance — track every LLM call"),
    ("Nexus", "Shared knowledge graphs as collective memory"),
    ("Prometheus", "Skill acquisition mid-session"),
    ("Crucible", "Quality benchmarking and red-teaming"),
    ("Agora", "Human-in-the-loop as equal participant"),
]
for i, (name, desc) in enumerate(caps):
    y = 2.3 + i * 0.65
    add_text(s, 0.8, y, 2.5, 0.4, name, size=16, color=ACCENT, bold=True)
    add_text(s, 3.5, y, 9, 0.4, desc, size=15, color=DIM)
slide_number(s, 8)

# ─── Slide 9: 13 Products ────────────────────────────────────────────
s = add_slide()
section_header(s, "Products", "13 products, one API")
left_products = [
    "Content Review Panels", "Living Personas", "Campaign Colosseum",
    "War Room", "Trend Forge", "Echo Chamber", "Consensus Engine",
]
right_products = [
    "Sentinel", "Replay", "Symposium (synthetic focus groups)",
    "Arena (A/B prediction)", "Archetype (persona generation)",
    "Augur (sentiment prediction)",
]
for i, p in enumerate(left_products):
    add_text(s, 0.8, 2.3 + i * 0.6, 5, 0.4, p, size=16, color=FG)
for i, p in enumerate(right_products):
    add_text(s, 6.8, 2.3 + i * 0.6, 6, 0.4, p, size=16, color=FG)
slide_number(s, 9)

# ─── Slide 10: Market Opportunity ────────────────────────────────────
s = add_slide()
section_header(s, "Market", "The opportunity")
add_text(s, 0.8, 2.3, 4, 0.8, "$80B+", size=48, color=ACCENT, bold=True)
add_text(s, 0.8, 3.2, 4, 0.4, "market research industry", size=16, color=DIM)
comparisons = [
    ("Focus group cost", "$10K-$50K", "~$100"),
    ("Time to insight", "2 weeks", "2 hours"),
    ("Segments tested", "1 at a time", "50 simultaneously"),
    ("Sentiment analysis", "Backward-looking", "Predictive"),
]
for i, (label, before, after) in enumerate(comparisons):
    y = 4.2 + i * 0.7
    add_text(s, 0.8, y, 3.5, 0.4, label, size=14, color=DIM)
    add_text(s, 4.5, y, 2.5, 0.4, before, size=14, color=DIM)
    add_text(s, 7.2, y, 0.5, 0.4, "->", size=14, color=ACCENT)
    add_text(s, 8.0, y, 3.5, 0.4, after, size=14, color=WHITE, bold=True)
slide_number(s, 10)

# ─── Slide 11: What's Built ──────────────────────────────────────────
s = add_slide()
section_header(s, "Progress", "What's built today")
add_bullet_list(s, 0.8, 2.2, 11, 5.0, [
    "Convex backend — fully deployed, live API, serverless",
    "Content review endpoint — submit, deliberate, verdict, follow-up questions",
    "10 reviewer personas with multi-round convergence tracking",
    "Vision model support — agents can see and analyze images",
    "Webhook callbacks — IMAI integration working end-to-end",
    "Landing page — Next.js, live at localhost:3100",
    "Open source — Apache 2.0, github.com/ZenderGoD/atherum",
    "Self-tested — Atherum reviewed itself, scored 88/100",
], size=17, color=FG)
slide_number(s, 11)

# ─── Slide 12: Architecture ──────────────────────────────────────────
s = add_slide()
section_header(s, "Architecture", "Fully serverless on Convex")
add_text(s, 0.8, 2.5, 11, 0.5, "Client  -->  Convex HTTP Actions  -->  LLM (OpenRouter)  -->  Convex Database", size=20, color=ACCENT)
add_bullet_list(s, 0.8, 3.5, 11, 3.5, [
    "No server to manage — Convex handles API, database, and background actions",
    "Deliberation runs as a Node.js action with full LLM orchestration",
    "All state persisted — reviews, rounds, agent responses, results",
    "Follow-up questions read the full transcript from database",
    "Scales automatically — no infrastructure to provision",
    "OASIS Python worker (planned) for social simulation",
], size=17, color=FG)
slide_number(s, 12)

# ─── Slide 13: IMAI Integration ──────────────────────────────────────
s = add_slide()
section_header(s, "First Customer", "IMAI Studio integration")
add_bullet_list(s, 0.8, 2.2, 11, 5.0, [
    "IMAI is an AI-powered creative studio for product visuals and marketing",
    "Users generate images -> select assets -> click Review",
    "Atherum reviews content via API, returns approval score + verdict",
    "Score shown on asset tile, full verdict in a dialog",
    "Accept / Reject / Revise workflow with comments and tags",
    "Review Inbox page for managing all reviews",
    "10 agents deliberate on each piece of content",
], size=17, color=FG)
slide_number(s, 13)

# ─── Slide 14: Tech Stack ────────────────────────────────────────────
s = add_slide()
section_header(s, "Technology", "Stack")
stack = [
    ("Backend", "Convex — serverless database + functions + API"),
    ("LLM", "OpenRouter — any model (GPT, Gemini, Claude, Nemotron)"),
    ("Frontend", "Next.js 15, Tailwind CSS v4"),
    ("Build", "pnpm workspaces + Turborepo"),
    ("Simulation", "OASIS / CAMEL-AI (Python)"),
    ("License", "Apache 2.0 (open source, no copyleft)"),
]
for i, (label, value) in enumerate(stack):
    y = 2.5 + i * 0.75
    add_text(s, 0.8, y, 2.5, 0.4, label, size=16, color=ACCENT, bold=True)
    add_text(s, 3.5, y, 9, 0.4, value, size=16, color=FG)
slide_number(s, 14)

# ─── Slide 15: Roadmap ───────────────────────────────────────────────
s = add_slide()
section_header(s, "Roadmap", "What's next")
phases = [
    ("Phase 1", "DONE", "Core engine, content review, IMAI integration"),
    ("Phase 2", "NEXT", "Worker pools, multi-tenancy, batch reviews, API keys"),
    ("Phase 3", "", "Hydra (recursive agents), OASIS simulation, better rubrics"),
    ("Phase 4", "", "Living Personas, Symposium, Arena, Archetype"),
    ("Phase 5", "", "Enterprise — Sentinel, War Room, Augur"),
]
for i, (phase, status, desc) in enumerate(phases):
    y = 2.4 + i * 0.9
    add_text(s, 0.8, y, 1.5, 0.4, phase, size=16, color=WHITE, bold=True)
    if status:
        sc = ACCENT if status == "DONE" else WHITE
        add_text(s, 2.5, y, 1.0, 0.4, status, size=12, color=sc, bold=True)
    add_text(s, 3.8, y, 9, 0.4, desc, size=16, color=DIM)
slide_number(s, 15)

# ─── Slide 16: Open Source ────────────────────────────────────────────
s = add_slide()
section_header(s, "Strategy", "Why open source")
add_bullet_list(s, 0.8, 2.2, 11, 4.5, [
    "Apache 2.0 — permissive, enterprise-friendly, patent protection",
    "Builds trust and drives adoption in the developer community",
    "Contributors improve the platform, add persona archetypes, find bugs",
    "Open core model: engine is free, hosted platform + enterprise features are paid",
    "Competitive moat: network effects from persona data and deliberation quality",
    "Clean-room implementation — no copyleft (AGPL/GPL) dependencies",
], size=17, color=FG)
slide_number(s, 16)

# ─── Slide 17: CTA ───────────────────────────────────────────────────
s = add_slide()
add_text(s, 0.8, 2.5, 11, 1.0, "Atherum", size=54, color=WHITE, bold=True)
add_text(s, 0.8, 3.7, 11, 0.6, "Collective Intelligence Engine", size=24, color=ACCENT)
add_text(s, 0.8, 5.0, 11, 0.4, "github.com/ZenderGoD/atherum", size=16, color=FG)
add_text(s, 0.8, 5.5, 11, 0.4, "API: https://next-okapi-818.convex.site", size=16, color=FG)
add_text(s, 0.8, 6.0, 11, 0.4, "Part of the Zeus Ecosystem", size=14, color=DIM)
slide_number(s, 17)

# Save
output = "/Users/bishalbanerjee/zeus/atherum/docs/Atherum-Stakeholder-Deck.pptx"
prs.save(output)
print(f"Saved: {output}")
